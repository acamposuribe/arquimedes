from __future__ import annotations

import json
from pathlib import Path

import pytest

from arquimedes.extract_office import (
    XLSX_MAX_ROWS_PER_SHEET,
    extract_raw_docx,
    extract_raw_pptx,
    extract_raw_xlsx,
)


def _read_jsonl(path: Path) -> list[dict]:
    return [json.loads(line) for line in path.read_text().splitlines() if line.strip()]


def _manifest(rel: str) -> dict:
    return {
        "file_hash": "deadbeef",
        "relative_path": rel,
        "domain": "research",
        "collection": "notes",
        "ingested_at": "2026-05-02T00:00:00+00:00",
    }


def test_extract_raw_docx(tmp_path):
    docx = pytest.importorskip("docx")
    src = tmp_path / "doc.docx"
    document = docx.Document()
    document.add_heading("Introduction", level=1)
    document.add_paragraph("First paragraph of the body.")
    document.add_heading("Methods", level=2)
    document.add_paragraph("Second paragraph here.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "header A"
    table.cell(0, 1).text = "header | B"
    table.cell(1, 0).text = "row a"
    table.cell(1, 1).text = "row | b"
    document.save(str(src))

    out = tmp_path / "extracted" / "mid"
    out.mkdir(parents=True)
    extract_raw_docx(src, out, "mid", _manifest("Research/notes/doc.docx"))

    pages = _read_jsonl(out / "pages.jsonl")
    chunks = _read_jsonl(out / "chunks.jsonl")
    text_md = (out / "text.md").read_text()
    meta = json.loads((out / "meta.json").read_text())

    assert "Introduction" in text_md
    assert "| header A | header \\| B |" in text_md
    assert "row \\| b" in text_md
    assert any("Introduction" in p["text"] for p in pages)
    assert chunks
    assert meta["file_type"] == "docx"
    assert meta["title"] == "Introduction"
    assert meta["page_count"] == len(pages)


def test_extract_raw_pptx(tmp_path):
    pptx_mod = pytest.importorskip("pptx")
    from pptx.util import Inches

    src = tmp_path / "deck.pptx"
    pres = pptx_mod.Presentation()
    title_only = pres.slide_layouts[5]
    blank = pres.slide_layouts[6]

    s1 = pres.slides.add_slide(title_only)
    s1.shapes.title.text = "Design Options"
    box = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Option A vs Option B"

    s2 = pres.slides.add_slide(blank)
    box2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box2.text_frame.text = "Second slide body"

    pres.save(str(src))

    out = tmp_path / "extracted" / "mid"
    out.mkdir(parents=True)
    extract_raw_pptx(src, out, "mid", _manifest("Practice/decks/deck.pptx"))

    pages = _read_jsonl(out / "pages.jsonl")
    meta = json.loads((out / "meta.json").read_text())

    assert len(pages) == 2
    assert pages[0]["headings"][0].startswith("Slide 1")
    assert "Design Options" in pages[0]["headings"][0]
    assert "Option A vs Option B" in pages[0]["text"]
    assert pages[1]["headings"][0].startswith("Slide 2")
    assert "Second slide body" in pages[1]["text"]
    assert meta["file_type"] == "pptx"
    assert meta["page_count"] == 2


def test_extract_raw_xlsx(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    src = tmp_path / "book.xlsx"
    wb = openpyxl.Workbook()
    s1 = wb.active
    s1.title = "Budget"
    s1.append(["Item", "Cost"])
    s1.append(["Rent", 1200])
    s1.append(["Food", 300])
    s2 = wb.create_sheet("Notes")
    s2.append(["Heading"])
    s2.append(["A note"])
    wb.save(str(src))

    out = tmp_path / "extracted" / "mid"
    out.mkdir(parents=True)
    extract_raw_xlsx(src, out, "mid", _manifest("Practice/sheets/book.xlsx"))

    pages = _read_jsonl(out / "pages.jsonl")
    meta = json.loads((out / "meta.json").read_text())

    assert len(pages) == 2
    assert pages[0]["headings"][0] == "Sheet: Budget"
    assert "Rent" in pages[0]["text"]
    assert "1200" in pages[0]["text"]
    assert pages[1]["headings"][0] == "Sheet: Notes"
    assert meta["file_type"] == "xlsx"
    assert meta["page_count"] == 2
    assert meta["title"] == "Budget"


def test_extract_raw_xlsx_fallback_without_expat(tmp_path, monkeypatch):
    openpyxl = pytest.importorskip("openpyxl")
    src = tmp_path / "fallback.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Item", "Cost"])
    ws.append(["Door", 1200])
    wb.save(src)

    def fail_load_workbook(*args, **kwargs):
        raise ImportError("No module named pyexpat")

    monkeypatch.setattr(openpyxl, "load_workbook", fail_load_workbook)
    out = tmp_path / "out-fallback"
    extract_raw_xlsx(src, out, "mid", _manifest("Practice/sheets/fallback.xlsx"))

    text = (out / "text.md").read_text(encoding="utf-8")
    assert "Sheet: Budget" in text
    assert "Door" in text
    assert "1200" in text
    meta = json.loads((out / "meta.json").read_text(encoding="utf-8"))
    assert any("limited XLSX fallback" in warning for warning in meta["extraction_warnings"])


def test_extract_raw_xlsx_truncates_large_sheets(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    src = tmp_path / "big.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Big"
    ws.append(["Idx", "Val"])
    for i in range(XLSX_MAX_ROWS_PER_SHEET + 50):
        ws.append([i, f"v{i}"])
    wb.save(str(src))

    out = tmp_path / "extracted" / "mid"
    out.mkdir(parents=True)
    extract_raw_xlsx(src, out, "mid", _manifest("Practice/sheets/big.xlsx"))

    warnings = _read_jsonl(out / "extraction_warnings.jsonl")
    assert any("truncated" in w["message"] for w in warnings)
