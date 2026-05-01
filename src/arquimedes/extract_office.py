"""Deterministic extraction for OOXML Office formats: .docx, .pptx, .xlsx."""

from __future__ import annotations

from pathlib import Path

from arquimedes.extract_text import (
    DEFAULT_PAGE_TARGET_CHARS,
    _pack_paragraphs_into_pages,
    write_synthetic_extraction,
)
from arquimedes.models import MaterialMeta, Page

# Bounds to keep large workbooks/decks from exhausting memory or page budgets.
XLSX_MAX_ROWS_PER_SHEET = 2000
XLSX_MAX_COLS_PER_SHEET = 64


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def _docx_paragraph_to_markdown(para) -> str:
    text = (para.text or "").strip()
    if not text:
        return ""
    style = (para.style.name if para.style is not None else "") or ""
    style_lower = style.lower()
    if style_lower.startswith("heading"):
        # "Heading 1" -> 1, default to 2 if unparseable
        level = 2
        digits = "".join(ch for ch in style_lower if ch.isdigit())
        if digits:
            level = max(1, min(6, int(digits)))
        return ("#" * level) + " " + text
    if style_lower.startswith("title"):
        return "# " + text
    if style_lower.startswith("list"):
        return "- " + text
    return text


def _markdown_table_cell(value: object) -> str:
    """Render a value safely inside a Markdown table cell."""
    return str(value or "").replace("\n", " ").replace("|", "\\|").strip()


def _docx_table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [_markdown_table_cell(cell.text) for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def extract_raw_docx(
    source_path: Path,
    output_dir: Path,
    material_id: str,
    manifest_entry: dict,
    chunk_size: int = 500,
) -> MaterialMeta:
    """Extract a .docx file: paragraphs + tables in document order, headings preserved."""
    from docx import Document
    from docx.oxml.ns import qn

    warnings_: list[str] = []
    title = ""

    try:
        doc = Document(str(source_path))
    except Exception as exc:
        warnings_.append(f"failed to open docx: {exc}")
        return write_synthetic_extraction(
            output_dir,
            material_id,
            manifest_entry,
            pages=[],
            text_md="",
            file_type="docx",
            chunk_size=chunk_size,
            warnings_=warnings_,
            classify=False,
        )

    # Iterate body elements in document order to keep paragraphs and tables interleaved.
    body = doc.element.body
    para_elements = {p._p: p for p in doc.paragraphs}
    table_elements = {t._tbl: t for t in doc.tables}

    blocks: list[str] = []
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            para = para_elements.get(child)
            if para is None:
                continue
            md = _docx_paragraph_to_markdown(para)
            if md:
                blocks.append(md)
                if not title and md.startswith("# "):
                    title = md[2:].strip()
        elif child.tag == qn("w:tbl"):
            table = table_elements.get(child)
            if table is None:
                continue
            md = _docx_table_to_markdown(table)
            if md:
                blocks.append(md)

    if not blocks:
        warnings_.append("docx contained no extractable paragraphs or tables")

    page_texts = _pack_paragraphs_into_pages(blocks) or [""]
    pages = [Page(page_number=i, text=t) for i, t in enumerate(page_texts, start=1) if t]
    if not pages:
        pages = [Page(page_number=1, text="")]

    text_md = "\n\n".join(blocks).strip() + ("\n" if blocks else "")

    return write_synthetic_extraction(
        output_dir,
        material_id,
        manifest_entry,
        pages,
        text_md,
        file_type="docx",
        title=title,
        chunk_size=chunk_size,
        warnings_=warnings_,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def _pptx_shape_text(shape) -> str:
    """Best-effort text extraction from a pptx shape."""
    pieces: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text or "" for run in para.runs).strip()
            if not line and para.text:
                line = para.text.strip()
            if line:
                pieces.append(line)
    if getattr(shape, "has_table", False):
        try:
            rows = []
            for row in shape.table.rows:
                cells = [_markdown_table_cell(cell.text) for cell in row.cells]
                rows.append(cells)
            if rows:
                width = max(len(r) for r in rows)
                rows = [r + [""] * (width - len(r)) for r in rows]
                header, body = rows[0], rows[1:]
                lines = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(["---"] * width) + " |",
                ]
                for r in body:
                    lines.append("| " + " | ".join(r) + " |")
                pieces.append("\n".join(lines))
        except Exception:
            pass
    return "\n".join(pieces).strip()


def extract_raw_pptx(
    source_path: Path,
    output_dir: Path,
    material_id: str,
    manifest_entry: dict,
    chunk_size: int = 500,
) -> MaterialMeta:
    """Extract a .pptx file: one synthetic page per slide."""
    from pptx import Presentation

    warnings_: list[str] = []
    title = ""

    try:
        pres = Presentation(str(source_path))
    except Exception as exc:
        warnings_.append(f"failed to open pptx: {exc}")
        return write_synthetic_extraction(
            output_dir,
            material_id,
            manifest_entry,
            pages=[],
            text_md="",
            file_type="pptx",
            chunk_size=chunk_size,
            warnings_=warnings_,
            classify=False,
        )

    pages: list[Page] = []
    text_md_parts: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        slide_title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                slide_title = (slide.shapes.title.text or "").strip()
        except Exception:
            slide_title = ""

        sorted_shapes = sorted(
            slide.shapes,
            key=lambda s: (
                getattr(s, "top", 0) or 0,
                getattr(s, "left", 0) or 0,
            ),
        )

        body_pieces: list[str] = []
        for shape in sorted_shapes:
            if shape == getattr(slide.shapes, "title", None):
                continue
            txt = _pptx_shape_text(shape)
            if txt:
                body_pieces.append(txt)

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    notes_text = (notes_frame.text or "").strip()
        except Exception:
            notes_text = ""

        heading_label = f"Slide {i}" + (f": {slide_title}" if slide_title else "")
        if not title and slide_title:
            title = slide_title

        page_md_parts = [f"# {heading_label}"]
        page_md_parts.extend(body_pieces)
        if notes_text:
            page_md_parts.append(f"## Speaker notes\n\n{notes_text}")
        page_text = "\n\n".join(page_md_parts).strip()

        page = Page(page_number=i, text=page_text)
        page.headings = [heading_label]
        pages.append(page)
        text_md_parts.append(page_text)

    if not pages:
        warnings_.append("pptx had no slides")
        pages = [Page(page_number=1, text="")]

    text_md = "\n\n".join(text_md_parts).strip() + ("\n" if text_md_parts else "")

    return write_synthetic_extraction(
        output_dir,
        material_id,
        manifest_entry,
        pages,
        text_md,
        file_type="pptx",
        title=title,
        chunk_size=chunk_size,
        warnings_=warnings_,
    )


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def _xlsx_rows_to_markdown(
    rows: list[list[str]],
    *,
    truncated_rows: bool,
    truncated_cols: bool,
) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    if truncated_rows:
        lines.append(f"| _truncated after {XLSX_MAX_ROWS_PER_SHEET} rows_ |" + " |" * (width - 1))
    if truncated_cols:
        lines.append(f"| _truncated after {XLSX_MAX_COLS_PER_SHEET} columns_ |" + " |" * (width - 1))
    return "\n".join(lines)


def extract_raw_xlsx(
    source_path: Path,
    output_dir: Path,
    material_id: str,
    manifest_entry: dict,
    chunk_size: int = 500,
) -> MaterialMeta:
    """Extract a .xlsx file: one synthetic page per visible worksheet."""
    from openpyxl import load_workbook

    warnings_: list[str] = []
    title = ""

    try:
        wb = load_workbook(filename=str(source_path), read_only=True, data_only=True)
    except Exception as exc:
        warnings_.append(f"failed to open xlsx: {exc}")
        return write_synthetic_extraction(
            output_dir,
            material_id,
            manifest_entry,
            pages=[],
            text_md="",
            file_type="xlsx",
            chunk_size=chunk_size,
            warnings_=warnings_,
            classify=False,
        )

    pages: list[Page] = []
    text_md_parts: list[str] = []
    page_index = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if getattr(ws, "sheet_state", "visible") != "visible":
            continue
        page_index += 1

        truncated_rows = False
        truncated_cols = False
        rows: list[list[str]] = []
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            if row_idx > XLSX_MAX_ROWS_PER_SHEET:
                truncated_rows = True
                break
            cells = list(row)
            if len(cells) > XLSX_MAX_COLS_PER_SHEET:
                truncated_cols = True
                cells = cells[:XLSX_MAX_COLS_PER_SHEET]
            rendered = [_markdown_table_cell(v) for v in cells]
            if any(c for c in rendered):
                rows.append(rendered)

        heading_label = f"Sheet: {sheet_name}"
        if not title:
            title = sheet_name

        if rows:
            table_md = _xlsx_rows_to_markdown(
                rows,
                truncated_rows=truncated_rows,
                truncated_cols=truncated_cols,
            )
            page_text = f"# {heading_label}\n\n{table_md}".strip()
        else:
            warnings_.append(f"sheet {sheet_name!r} was empty")
            page_text = f"# {heading_label}\n\n_empty sheet_"

        if truncated_rows:
            warnings_.append(
                f"sheet {sheet_name!r} truncated at {XLSX_MAX_ROWS_PER_SHEET} rows"
            )
        if truncated_cols:
            warnings_.append(
                f"sheet {sheet_name!r} truncated at {XLSX_MAX_COLS_PER_SHEET} columns"
            )

        page = Page(page_number=page_index, text=page_text)
        page.headings = [heading_label]
        pages.append(page)
        text_md_parts.append(page_text)

    wb.close()

    if not pages:
        warnings_.append("xlsx had no visible worksheets")
        pages = [Page(page_number=1, text="")]

    text_md = "\n\n".join(text_md_parts).strip() + ("\n" if text_md_parts else "")

    return write_synthetic_extraction(
        output_dir,
        material_id,
        manifest_entry,
        pages,
        text_md,
        file_type="xlsx",
        title=title,
        chunk_size=chunk_size,
        warnings_=warnings_,
    )


__all__ = [
    "extract_raw_docx",
    "extract_raw_pptx",
    "extract_raw_xlsx",
]


# Silence unused import warning while keeping helper available for future use.
_ = DEFAULT_PAGE_TARGET_CHARS
